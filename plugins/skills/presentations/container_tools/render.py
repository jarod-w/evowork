#!/usr/bin/env python3
"""内容 JSON → pptx 渲染器（08 §5.3 的结构化生成）。

模型的输出是内容与结构，**排版由这里决定**。这么切的理由写在 SKILL.md 里，
这里只补一条实现层面的：所有字号、留白、配色都来自 `templates/<id>/template.json`，
渲染代码里不出现字面量数值 —— 否则"换模板"就会变成"改渲染代码"。

退出码、校验失败的措辞、以及"报错里不许出现用户内容"这条纪律都在
`plugins/skills/_shared/evowork_skill.py` 里 —— 四个技能必须一致，见那个文件的头注释。
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

SKILL_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(SKILL_ROOT.parent / "_shared"))

from evowork_skill import (  # noqa: E402 -- 必须在 sys.path 调整之后
    ensure_office_runtime,
    EXIT_INVALID_CONTENT,
    EXIT_MISSING_ASSET,
    EXIT_OK,
    EXIT_RUNTIME_MISSING,
    fail,
    load_template,
    read_content,
    resolve_asset as _resolve_asset,
    runtime_missing_message,
    succeeded,
    validate_content,
)

LAYOUTS = ("title", "bullets", "chart", "table", "section")


def validate(content: dict) -> None:
    """先校验再渲染。具体的报错措辞与"不含用户内容"的纪律在共用骨架里。"""
    validate_content(
        content,
        SKILL_ROOT / "schema" / "content.schema.json",
        discriminator="layout",
        known_kinds=LAYOUTS,
    )


def resolve_asset(ref: str, base: Path) -> Path:
    return _resolve_asset(ref, base, hint="先用 charts 技能生成它，再引用。")


def render(content: dict, out_path: Path, base_dir: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu, Pt
    except ModuleNotFoundError:
        # 08 §4：办公扩展按需下载。文案要与解析管道**统一**——
        # 不能一处说"解析组件"一处说"生成组件"，那会让用户以为要装两个东西
        # 文案由共用骨架统一给（08 §4：不能一处说"解析组件"一处说"生成组件"）
        fail(EXIT_RUNTIME_MISSING, runtime_missing_message("office", "生成 pptx"))

    template = load_template(SKILL_ROOT, content.get("template", "business"))
    theme = template["theme"]
    metrics = template["metrics"]

    def color(name: str) -> "RGBColor":
        return RGBColor.from_string(theme[name].lstrip("#"))

    prs = Presentation()
    prs.slide_width = Emu(int(metrics["slide_width_emu"]))
    prs.slide_height = Emu(int(metrics["slide_height_emu"]))
    blank = prs.slide_layouts[6]  # 空白版式：所有元素由我们放，避免继承默认主题的字体

    def add_textbox(slide, spec_key: str, text: str, *, style: str) -> None:
        box = metrics[spec_key]
        shape = slide.shapes.add_textbox(
            Emu(int(box["left"])), Emu(int(box["top"])), Emu(int(box["width"])), Emu(int(box["height"]))
        )
        frame = shape.text_frame
        frame.word_wrap = True
        frame.text = text
        style_spec = theme["styles"][style]
        for paragraph in frame.paragraphs:
            paragraph.font.size = Pt(style_spec["size_pt"])
            paragraph.font.bold = style_spec.get("bold", False)
            paragraph.font.color.rgb = color(style_spec["color"])
            # 中文字体必须显式设置，否则在没有装对应字体的机器上会变方框
            paragraph.font.name = theme["font_cjk"]

    for spec in content["slides"]:
        layout = spec["layout"]
        slide = prs.slides.add_slide(blank)

        if layout == "title":
            add_textbox(slide, "title_slide_title", spec["title"], style="display")
            if spec.get("subtitle"):
                add_textbox(slide, "title_slide_subtitle", spec["subtitle"], style="subtitle")

        elif layout == "section":
            add_textbox(slide, "section_title", spec["title"], style="section")

        elif layout == "bullets":
            add_textbox(slide, "body_title", spec["title"], style="title")
            box = metrics["body_content"]
            shape = slide.shapes.add_textbox(
                Emu(int(box["left"])), Emu(int(box["top"])), Emu(int(box["width"])), Emu(int(box["height"]))
            )
            frame = shape.text_frame
            frame.word_wrap = True
            style_spec = theme["styles"]["body"]
            # 条目多时按模板给的梯度缩字号，而不是让文字溢出框
            size = style_spec["size_pt"]
            for step in metrics["bullet_shrink_steps"]:
                if len(spec["bullets"]) >= step["min_items"]:
                    size = step["size_pt"]
            for index, bullet in enumerate(spec["bullets"]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = f"· {bullet}"
                paragraph.font.size = Pt(size)
                paragraph.font.color.rgb = color(style_spec["color"])
                paragraph.font.name = theme["font_cjk"]

        elif layout == "chart":
            add_textbox(slide, "body_title", spec["title"], style="title")
            image = resolve_asset(spec["chart_ref"], base_dir)
            box = metrics["chart_area"]
            slide.shapes.add_picture(
                str(image), Emu(int(box["left"])), Emu(int(box["top"])), height=Emu(int(box["height"]))
            )
            if spec.get("caption"):
                add_textbox(slide, "caption", spec["caption"], style="caption")

        elif layout == "table":
            add_textbox(slide, "body_title", spec["title"], style="title")
            table_spec = spec["table"]
            rows = len(table_spec["rows"]) + 1
            cols = len(table_spec["header"])
            box = metrics["table_area"]
            shape = slide.shapes.add_table(
                rows, cols, Emu(int(box["left"])), Emu(int(box["top"])), Emu(int(box["width"])), Emu(int(box["height"]))
            )
            table = shape.table
            for col, name in enumerate(table_spec["header"]):
                cell = table.cell(0, col)
                cell.text = name
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(theme["styles"]["table_header"]["size_pt"])
                    paragraph.font.bold = True
                    paragraph.font.name = theme["font_cjk"]
            for row_index, row in enumerate(table_spec["rows"], start=1):
                for col in range(cols):
                    cell = table.cell(row_index, col)
                    cell.text = row[col] if col < len(row) else ""
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(theme["styles"]["table_body"]["size_pt"])
                        paragraph.font.name = theme["font_cjk"]
            if spec.get("caption"):
                add_textbox(slide, "caption", spec["caption"], style="caption")

        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    parser = argparse.ArgumentParser(description="把内容 JSON 渲染成 pptx")
    parser.add_argument("--content", required=True, help="内容 JSON 路径")
    parser.add_argument("--out", required=True, help="输出 pptx 路径")
    parser.add_argument(
        "--validate-only",
        action="store_true",
        help="只校验不渲染（不需要办公扩展，用于让模型先把内容改对）",
    )
    args = parser.parse_args()

    # 缺办公扩展的模块时换到扩展的解释器重跑（见 evowork_skill.ensure_office_runtime）。
    # --validate-only 不需要它，但提前换掉更简单，也让两条路径的行为一致
    ensure_office_runtime(("pptx",))

    content_path = Path(args.content).resolve()
    content = read_content(content_path)
    validate(content)

    if args.validate_only:
        succeeded(slides=len(content["slides"]), validated=True)
        return

    out_path = Path(args.out).resolve()
    render(content, out_path, content_path.parent)
    succeeded(out_path, slides=len(content["slides"]))


if __name__ == "__main__":
    main()
